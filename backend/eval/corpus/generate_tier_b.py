"""Deterministic generator for tier-B (synthetic edge-case documents).

Tier B exists to exercise every extractor path and every known boundary
condition — including ones no real corpus reliably contains (an OCR-only page,
a single CSV row wider than the table char budget). Because these documents are
authored here, their provenance is unambiguous: licence is the repository's, and
`source_url` is this generator.

Reproducibility — MEASURED, not assumed (re-run the script twice and diff):

  byte-identical across runs : .txt .md .csv .docx .pptx
                               (OOXML zips are rewritten with fixed entry
                               timestamps by _norm_zip)
  NOT byte-identical         : .xlsx .pdf

  openpyxl and PyMuPDF both embed time-varying state that survives the fixes
  attempted here (pinned docProps timestamps, pinned /CreationDate, /ModDate
  and /ID). The residual source was not chased further because byte-stable
  regeneration is NOT what protects a baseline.

  What protects a baseline is that the COMMITTED bytes are pinned: every file
  is checksummed in MANIFEST.yaml and verify.py hard-fails before any eval run
  if a single byte differs. Treat this generator as an audit and authoring aid
  — regenerating .xlsx/.pdf produces equivalent CONTENT with different bytes,
  which will (correctly) fail verify.py until the manifest is updated
  deliberately.

Usage (inside the api container, which has every parser installed):
    python eval/corpus/generate_tier_b.py            # write files
    python eval/corpus/generate_tier_b.py --manifest # emit YAML entries
"""
from __future__ import annotations

import argparse
import hashlib
import datetime
import shutil
import zipfile
from pathlib import Path

OUT = Path(__file__).resolve().parent / "tier-b"
FIXED_ZIP_DATE = (2020, 1, 1, 0, 0, 0)  # any constant; only stability matters
RETRIEVED = "2026-07-14"

# path -> (language, parallel_id, exercises)
SPEC: dict[str, tuple[str, str | None, list[str]]] = {}


def _norm_zip(path: Path) -> None:
    """Rewrite an OOXML archive with fixed entry timestamps + stored order."""
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for name in sorted(src.namelist()):
            info = zipfile.ZipInfo(name, date_time=FIXED_ZIP_DATE)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            dst.writestr(info, src.read(name))
    shutil.move(str(tmp), str(path))


# --- generators -------------------------------------------------------------

def gen_txt() -> None:
    (OUT / "b01-policy-notes.txt").write_text(
        "Retention Policy Notes\n"
        "Records are retained for seven years from the date of creation.\n"
        "Disposal requires written authorisation from the records officer.\n",
        encoding="utf-8", newline="\n",
    )
    SPEC["tier-b/b01-policy-notes.txt"] = ("en", None, ["extractor:txt", "plain-text"])


def gen_md() -> None:
    (OUT / "b02-handbook.md").write_text(
        "# Field Handbook\n\n"
        "## Reporting\n\n"
        "Incidents must be reported within 24 hours of discovery.\n\n"
        "## Escalation\n\n"
        "Unresolved incidents escalate to the duty supervisor after 72 hours.\n",
        encoding="utf-8", newline="\n",
    )
    SPEC["tier-b/b02-handbook.md"] = (
        "en", None, ["extractor:txt", "heading-path", "markdown-headings"],
    )


def gen_csv_normal() -> None:
    rows = ["region,quarter,units_sold,revenue_cad"]
    rows += [
        f"{r},{q},{100 + i * 7},{(100 + i * 7) * 12}"
        for i, (r, q) in enumerate(
            [(r, q) for r in ("North", "South", "East", "West") for q in ("Q1", "Q2", "Q3", "Q4")]
        )
    ]
    (OUT / "b03-sales.csv").write_text("\n".join(rows) + "\n", encoding="utf-8", newline="\n")
    SPEC["tier-b/b03-sales.csv"] = (
        "en", None, ["extractor:csv", "table-answer", "row-group-chunking", "header-repetition"],
    )


def gen_csv_oversized_row() -> None:
    """A single row far wider than table_chunk_char_budget (1600 chars).

    dynamic_rows_per_chunk() clamps with max(1, ...) at helper.py:18-21, so this
    row bypasses the budget entirely and produces one oversized chunk. With the
    embedding window at 2048 tokens (finding #19), the tail of this row is
    silently dropped from the dense vector — this document is what makes that
    failure observable in eval instead of theoretical.
    """
    filler = " ".join(f"clause{i:04d}" for i in range(1200))  # ~13k chars
    (OUT / "b04-wide-row.csv").write_text(
        "contract_id,summary\n"
        f"CT-9001,\"The indemnity ceiling is 4200000 CAD. {filler} "
        "The termination notice period is 95 days.\"\n",
        encoding="utf-8", newline="\n",
    )
    SPEC["tier-b/b04-wide-row.csv"] = (
        "en", None,
        ["extractor:csv", "long-table-tail", "oversized-single-row",
         "bypasses-helper-18-21", "embedding-window-2048"],
    )


def gen_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    # openpyxl stamps docProps/core.xml with the current time, which changes
    # the archive CONTENT (not just the zip mtime), so normalising the zip is
    # not enough on its own.
    fixed = datetime.datetime(2020, 1, 1)
    wb.properties.created = fixed
    wb.properties.modified = fixed
    wb.properties.creator = "knowall-eval-corpus"
    wb.properties.lastModifiedBy = "knowall-eval-corpus"
    ws = wb.active
    ws.title = "Inventory"
    ws.append(["sku", "warehouse", "on_hand"])
    for i in range(12):
        ws.append([f"SKU-{i:03d}", "Halifax" if i % 2 else "Regina", 40 + i])
    ws2 = wb.create_sheet("Thresholds")
    ws2.append(["metric", "value"])
    ws2.append(["reorder_point", 25])
    ws2.append(["safety_stock", 15])
    path = OUT / "b05-inventory.xlsx"
    wb.save(path)
    _norm_zip(path)
    SPEC["tier-b/b05-inventory.xlsx"] = (
        "en", None, ["extractor:xlsx", "table-answer", "multi-sheet", "sheet-context"],
    )


def gen_docx() -> None:
    from docx import Document as Docx

    doc = Docx()
    doc.add_heading("Operations Manual", level=1)
    doc.add_heading("Access Control", level=2)
    doc.add_paragraph(
        "Badge access is revoked automatically after 30 days of inactivity."
    )
    doc.add_heading("Maintenance Windows", level=2)
    doc.add_paragraph("Scheduled maintenance occurs on the second Sunday monthly.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "system"
    table.rows[0].cells[1].text = "window_hours"
    for name, hours in (("billing", "4"), ("reporting", "2")):
        cells = table.add_row().cells
        cells[0].text, cells[1].text = name, hours
    path = OUT / "b06-operations.docx"
    doc.save(path)
    _norm_zip(path)
    SPEC["tier-b/b06-operations.docx"] = (
        "en", None,
        ["extractor:docx", "heading-path", "docx-table", "in-position-tables"],
    )


def gen_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    for title, body in (
        ("Quarterly Review", "Customer churn fell to 3.1 percent."),
        ("Next Steps", "Migrate the reporting pipeline before the freeze."),
    ):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    path = OUT / "b07-review.pptx"
    prs.save(path)
    _norm_zip(path)
    SPEC["tier-b/b07-review.pptx"] = (
        "en", None, ["extractor:pptx", "slide-chunking"],
    )


def _pin_pdf(doc) -> None:
    """Pin every time-varying field PyMuPDF writes.

    An empty set_metadata() is not sufficient: PyMuPDF still emits
    /CreationDate, /ModDate and a content+time derived /ID, all of which
    change the bytes on every run.
    """
    stamp = "D:20200101000000Z"
    doc.set_metadata({
        "producer": "knowall-eval-corpus",
        "creator": "knowall-eval-corpus",
        "creationDate": stamp,
        "modDate": stamp,
        "title": "", "author": "", "subject": "", "keywords": "",
    })
    doc.xref_set_key(-1, "ID", "[<00000000000000000000000000000000>"
                               "<00000000000000000000000000000000>]")


def gen_pdf_text() -> None:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 96), "Service Level Agreement", fontsize=16)
    page.insert_text((72, 130), "Priority one tickets are acknowledged within 15 minutes.",
                     fontsize=11)
    page.insert_text((72, 150), "Credits apply when uptime falls below 99.5 percent.",
                     fontsize=11)
    _pin_pdf(doc)
    doc.save(OUT / "b08-sla.pdf", deflate=True, garbage=4)
    doc.close()
    SPEC["tier-b/b08-sla.pdf"] = (
        "en", None, ["extractor:pdf", "pdf-text-layer"],
    )


def gen_pdf_scanned() -> None:
    """Image-only PDF: no text layer, so extraction must fall through to OCR."""
    import pymupdf

    src = pymupdf.open()
    page = src.new_page()
    page.insert_text((72, 100), "ARCHIVED NOTICE", fontsize=20)
    page.insert_text((72, 140), "The heritage grant ceiling is 75000 dollars.", fontsize=14)
    pix = page.get_pixmap(dpi=200)
    src.close()

    doc = pymupdf.open()
    out_page = doc.new_page(width=pix.width * 72 / 200, height=pix.height * 72 / 200)
    out_page.insert_image(out_page.rect, pixmap=pix)
    _pin_pdf(doc)
    doc.save(OUT / "b09-scanned-notice.pdf", deflate=True, garbage=4)
    doc.close()
    SPEC["tier-b/b09-scanned-notice.pdf"] = (
        "en", None, ["extractor:pdf", "ocr-answer", "image-only-page", "no-text-layer"],
    )


def gen_distractors() -> None:
    """Topically adjacent, answer-free. These exist so a retriever that ignores
    the question still scores badly, and so abstention has something to reject."""
    specs = [
        ("b10-distractor-logistics.txt",
         "Logistics Overview\n"
         "Freight consolidation reduces handling costs across regional hubs.\n"
         "Carrier selection is reviewed annually by the procurement committee.\n"),
        ("b11-distractor-training.txt",
         "Training Catalogue\n"
         "Introductory workshops run monthly for new analysts.\n"
         "Advanced modules require completion of the introductory track.\n"),
        ("b12-distractor-facilities.txt",
         "Facilities Bulletin\n"
         "Lobby renovations continue through the end of the fiscal year.\n"
         "Parking permits are reissued each September.\n"),
    ]
    for name, text in specs:
        (OUT / name).write_text(text, encoding="utf-8", newline="\n")
        SPEC[f"tier-b/{name}"] = ("en", None, ["distractor", "no-answers"])


GENERATORS = (
    gen_txt, gen_md, gen_csv_normal, gen_csv_oversized_row, gen_xlsx,
    gen_docx, gen_pptx, gen_pdf_text, gen_pdf_scanned, gen_distractors,
)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(1 << 20), b""):
            digest.update(block)
    return digest.hexdigest()


def emit_manifest_entries() -> str:
    lines: list[str] = []
    for rel in sorted(SPEC):
        language, parallel_id, exercises = SPEC[rel]
        path = OUT.parent / rel
        lines += [
            f"  - path: {rel}",
            f"    sha256: {sha256_file(path)}",
            f"    format: {path.suffix.lstrip('.')}",
            f"    language: {language}",
            f"    parallel_id: {parallel_id if parallel_id else 'null'}",
            "    tier: b",
            "    license: CC0-1.0  # authored in this repository",
            "    source_url: generated by eval/corpus/generate_tier_b.py",
            f"    retrieved: {RETRIEVED}",
            "    exercises: [" + ", ".join(exercises) + "]",
        ]
    return "\n".join(lines)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--manifest", action="store_true",
                        help="print MANIFEST.yaml entries for the generated files")
    args = parser.parse_args()

    OUT.mkdir(parents=True, exist_ok=True)
    for generator in GENERATORS:
        generator()

    if args.manifest:
        print(emit_manifest_entries())
    else:
        for rel in sorted(SPEC):
            print(f"{rel}  {sha256_file(OUT.parent / rel)[:16]}…")
        print(f"\n{len(SPEC)} tier-B documents written to {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
