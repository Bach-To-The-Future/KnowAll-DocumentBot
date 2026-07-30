import logging
import os

from pptx import Presentation

from core.interfaces import ChunkLike
from extraction.base import BaseExtractor
from extraction.helper import generate_metadata

logging.basicConfig(level=logging.INFO)

class ExtractPPTX(BaseExtractor):
    def extract_and_chunk(self, file_path: str) -> list[ChunkLike]:
        """Extract and chunk PPTX content with metadata."""
        if not self.validate_file(file_path):
            return []

        self.logger.info(f"Extracting and chunking: {file_path}")
        source = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[-1][1:].lower()

        try:
            prs = Presentation(file_path)
            all_nodes: list[ChunkLike] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text = "\n".join(slide_text)
                    doc = self.create_document(text, metadata=generate_metadata(
                        source=source,
                        index=slide_num,
                        max_index=len(prs.slides),
                        file_format=ext,
                        page_number=slide_num,
                        content_type="text"
                    ))
                    if doc:
                        nodes = self.chunk_document(doc)
                        for i, node in enumerate(nodes):
                            node.metadata.update({"chunk_index": i, "total_chunks": len(nodes)})
                        all_nodes.extend(nodes)

            self.logger.info(f"Extracted {len(all_nodes)} chunks from {file_path}")
            return all_nodes
        except Exception as e:
            self.logger.error(f"Error processing PPTX file: {e}")
            return []
